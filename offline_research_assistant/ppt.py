from __future__ import annotations

from pathlib import Path
from typing import Dict, List, Sequence, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from .text_utils import split_sentences


THEMES: Dict[str, Dict[str, object]] = {
    "professional": {
        "title_font": "Calibri",
        "title_size": Pt(40),
        "title_color": RGBColor(31, 78, 121),
        "title_bold": True,
        "body_font": "Calibri",
        "body_size": Pt(18),
        "body_color": RGBColor(64, 64, 64),
        "background_color": RGBColor(255, 255, 255),
    },
    "modern": {
        "title_font": "Arial",
        "title_size": Pt(38),
        "title_color": RGBColor(47, 84, 150),
        "title_bold": True,
        "body_font": "Arial",
        "body_size": Pt(16),
        "body_color": RGBColor(89, 89, 89),
        "background_color": RGBColor(248, 248, 248),
    },
    "creative": {
        "title_font": "Georgia",
        "title_size": Pt(38),
        "title_color": RGBColor(192, 80, 77),
        "title_bold": True,
        "body_font": "Georgia",
        "body_size": Pt(18),
        "body_color": RGBColor(79, 79, 79),
        "background_color": RGBColor(252, 248, 227),
    },
}


def _shorten_for_bullet(text: str, max_chars: int = 180) -> str:
    t = (text or "").strip()
    if len(t) <= max_chars:
        return t
    cut = t[: max_chars - 1]
    for sep in ("; ", ". ", ", ", ": ", " - "):
        idx = cut.rfind(sep)
        if idx >= int(max_chars * 0.6):
            cut = cut[:idx]
            break
    return cut.rstrip(" ,;:-") + "…"


def _apply_slide_background(slide, background_color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = background_color


def _add_top_bar(slide, color: RGBColor) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0),
        top=Inches(0),
        width=Inches(10),
        height=Inches(0.32),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def _add_title(slide, title: str, theme: Dict[str, object]) -> None:
    _apply_slide_background(slide, theme["background_color"])  # type: ignore
    _add_top_bar(slide, theme["title_color"])  # type: ignore

    box = slide.shapes.add_textbox(
        left=Inches(0.7),
        top=Inches(0.55),
        width=Inches(8.9),
        height=Inches(0.8),
    )
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (title or "").strip() or "Slide"
    p.alignment = PP_ALIGN.LEFT
    p.font.name = theme["title_font"]  # type: ignore
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = theme["title_color"]  # type: ignore


def _content_box(slide, *, left: float, top: float, width: float, height: float):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))


def _normalize_bullets(lines: Sequence[str]) -> List[str]:
    out: List[str] = []
    for b in lines:
        t = (b or "").strip()
        if not t:
            continue
        # remove any leading bullet chars to avoid duplicates
        t = t.lstrip("-•* ")
        out.append(t)
    return out


def _chunk(items: Sequence[str], size: int) -> List[List[str]]:
    if size <= 0:
        return [list(items)]
    chunks: List[List[str]] = []
    cur: List[str] = []
    for it in items:
        cur.append(it)
        if len(cur) >= size:
            chunks.append(cur)
            cur = []
    if cur:
        chunks.append(cur)
    return chunks


def _add_bullets(
    slide,
    bullets: Sequence[str],
    theme: Dict[str, object],
    *,
    left: float = 0.95,
    top: float = 1.55,
    width: float = 11.4,
    height: float = 5.5,
    font_size: Pt | None = None,
) -> None:
    box = _content_box(slide, left=left, top=top, width=width, height=height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE

    bullets = _normalize_bullets(bullets)
    if font_size is None:
        font_size = theme.get("body_size", Pt(18))  # type: ignore
    if not bullets:
        bullets = ["(No content)"]

    added = 0
    for b in bullets:
        b = _shorten_for_bullet(b, max_chars=190)
        if not b:
            continue
        p = tf.paragraphs[0] if added == 0 else tf.add_paragraph()
        p.text = "• " + b
        p.level = 0
        p.space_after = Pt(6)
        p.line_spacing = 1.12
        p.font.name = theme["body_font"]  # type: ignore
        p.font.size = font_size
        p.font.color.rgb = theme["body_color"]  # type: ignore
        added += 1


def _add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.3))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER


def _bucket_sentences(summary: str) -> Dict[str, List[str]]:
    s = split_sentences(summary)
    if not s:
        return {"Context": [], "Method": [], "Results": [], "Limitations": [], "Conclusion": []}

    buckets = {"Context": [], "Method": [], "Results": [], "Limitations": [], "Conclusion": []}
    for sent in s:
        low = sent.lower()
        if any(k in low for k in ["method", "dataset", "experiment", "architecture", "model", "approach"]):
            buckets["Method"].append(sent)
        elif any(k in low for k in ["result", "accuracy", "improve", "outperform", "achieve", "metric"]):
            buckets["Results"].append(sent)
        elif any(k in low for k in ["limitation", "however", "constraint", "future work", "future"]):
            buckets["Limitations"].append(sent)
        elif any(k in low for k in ["conclude", "overall", "in summary", "we show", "this paper"]):
            buckets["Conclusion"].append(sent)
        else:
            buckets["Context"].append(sent)

    return buckets


def _derive_deck_sections(summary_sents: List[str]) -> Dict[str, List[str]]:
    buckets = {
        "Problem & Motivation": [],
        "Key Contributions": [],
        "Method Overview": [],
        "Experimental Setup": [],
        "Results & Findings": [],
        "Limitations": [],
        "Future Work": [],
        "Conclusion": [],
    }

    for s in summary_sents:
        low = s.lower()
        if any(k in low for k in ["we propose", "we present", "we introduce", "we develop", "our method", "our approach", "framework"]):
            buckets["Key Contributions"].append(s)
            buckets["Method Overview"].append(s)
        elif any(k in low for k in ["dataset", "datasets", "benchmark", "evaluation", "evaluate", "experiment", "experiments", "ablation", "baseline"]):
            buckets["Experimental Setup"].append(s)
        elif any(k in low for k in ["result", "results", "outperform", "improve", "improves", "accuracy", "f1", "auc", "precision", "recall", "error", "gain"]):
            buckets["Results & Findings"].append(s)
        elif any(k in low for k in ["limitation", "limitations", "however", "constraint", "trade-off", "threat", "threats"]):
            buckets["Limitations"].append(s)
        elif any(k in low for k in ["future work", "future", "next", "further", "extension", "extend"]):
            buckets["Future Work"].append(s)
        elif any(k in low for k in ["in summary", "overall", "we conclude", "conclude", "conclusion"]):
            buckets["Conclusion"].append(s)
        else:
            buckets["Problem & Motivation"].append(s)

    # De-duplicate while preserving order
    for k, vals in list(buckets.items()):
        seen = set()
        out: List[str] = []
        for v in vals:
            key = v.strip()
            if not key or key in seen:
                continue
            seen.add(key)
            out.append(v)
        buckets[k] = out
    return buckets


def _add_section_slides(
    prs: Presentation,
    title_base: str,
    bullets: Sequence[str],
    theme: Dict[str, object],
    *,
    max_bullets_per_slide: int = 6,
    font_size: Pt | None = None,
) -> int:
    bullets = _normalize_bullets(bullets)
    if not bullets:
        return 0
    pages = _chunk(list(bullets), max_bullets_per_slide)
    for idx, page in enumerate(pages, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = title_base if len(pages) == 1 else f"{title_base} ({idx}/{len(pages)})"
        _add_title(slide, title, theme)
        _add_bullets(slide, page, theme, font_size=font_size)
    return len(pages)


def create_ppt_from_summary(
    summary: str,
    output_path: str | Path,
    theme_name: str = "professional",
    title: str = "Research Summary",
    subtitle: str | None = None,
    keywords: Sequence[str] | None = None,
) -> Path:
    theme = THEMES.get(theme_name, THEMES["professional"])
    out_pptx = Path(output_path)
    out_pptx.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    # Widescreen 16:9 layout (better for modern decks)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    summary_sents = [s.strip() for s in split_sentences(summary) if s.strip()]
    deck = _derive_deck_sections(summary_sents)

    # Ensure key sections have content so the deck feels complete even when
    # the summary doesn't contain obvious cue words.
    if summary_sents:
        if not deck.get("Method Overview"):
            deck["Method Overview"] = summary_sents[1:5]
        if not deck.get("Results & Findings"):
            deck["Results & Findings"] = summary_sents[5:9]
        if not deck.get("Limitations"):
            deck["Limitations"] = [s for s in summary_sents[-6:] if s][:3]
        if not deck.get("Conclusion"):
            deck["Conclusion"] = [s for s in summary_sents[-3:] if s]

    # 1) Title slide (clean, centered)
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_slide_background(s1, theme["background_color"])  # type: ignore
    _add_top_bar(s1, theme["title_color"])  # type: ignore

    title_box = s1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.3))
    ttf = title_box.text_frame
    ttf.clear()
    ttf.word_wrap = True
    p = ttf.paragraphs[0]
    p.text = (title or "").strip() or "Research Summary"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = theme["title_font"]  # type: ignore
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = theme["title_color"]  # type: ignore

    sub_text = (subtitle or "Generated offline")
    sub_box = s1.shapes.add_textbox(Inches(1.1), Inches(3.35), Inches(11.1), Inches(0.7))
    stf = sub_box.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    sp.text = sub_text
    sp.alignment = PP_ALIGN.CENTER
    sp.font.name = theme["body_font"]  # type: ignore
    sp.font.size = Pt(22)
    sp.font.color.rgb = theme["body_color"]  # type: ignore

    # 2) Paper-at-a-glance (more useful than a generic agenda)
    s2a = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(s2a, "Paper at a Glance", theme)
    at_a_glance: List[str] = []
    if summary_sents:
        at_a_glance.append(summary_sents[0])
    if deck.get("Method Overview"):
        at_a_glance.append(deck["Method Overview"][0])
    if deck.get("Results & Findings"):
        at_a_glance.append(deck["Results & Findings"][0])
    if keywords:
        kw = [k.strip() for k in (keywords or []) if k and k.strip()]
        if kw:
            at_a_glance.append("Keywords: " + ", ".join(kw[:8]))
    _add_bullets(s2a, at_a_glance[:6], theme, font_size=Pt(18))

    # 3) Agenda
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(s2, "Agenda", theme)
    _add_bullets(
        s2,
        [
            "Problem & motivation",
            "Key contributions",
            "Method overview",
            "Experimental setup",
            "Results & findings",
            "Limitations and future work",
            "Conclusion and takeaways",
        ],
        theme,
    )

    # 4+) Executive summary (spans multiple slides if needed)
    if not summary_sents:
        _add_section_slides(prs, "Executive Summary", ["(No summary text provided)"] , theme)
    else:
        _add_section_slides(prs, "Executive Summary", summary_sents[:18], theme, max_bullets_per_slide=6)

    # Core sections (auto-paginated)
    _add_section_slides(prs, "Problem & Motivation", deck.get("Problem & Motivation", []), theme)
    _add_section_slides(prs, "Key Contributions", deck.get("Key Contributions", []), theme)
    _add_section_slides(prs, "Method Overview", deck.get("Method Overview", []), theme)
    _add_section_slides(prs, "Experimental Setup", deck.get("Experimental Setup", []), theme)
    _add_section_slides(prs, "Results & Findings", deck.get("Results & Findings", []), theme)
    _add_section_slides(prs, "Limitations", deck.get("Limitations", []), theme)
    _add_section_slides(prs, "Future Work", deck.get("Future Work", []), theme)
    _add_section_slides(prs, "Conclusion", deck.get("Conclusion", []), theme)

    # Detailed summary (if there is still remaining content)
    if len(summary_sents) > 18:
        _add_section_slides(prs, "Detailed Notes", summary_sents[18:], theme, max_bullets_per_slide=7, font_size=Pt(16))

    # Keywords / Glossary (optional)
    if keywords:
        kw = [k.strip() for k in keywords if k and k.strip()]
        if kw:
            s = prs.slides.add_slide(prs.slide_layouts[6])
            _add_title(s, "Keywords", theme)
            _add_bullets(s, kw[:18], theme, font_size=Pt(16))

    # Closing slide
    s_end = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(s_end, "Q&A", theme)
    _add_bullets(s_end, ["Questions", "Discussion"], theme)

    # Footer on all slides
    for slide in prs.slides:
        _add_footer(slide, "Offline Research Assistant")

    prs.save(str(out_pptx))
    return out_pptx
